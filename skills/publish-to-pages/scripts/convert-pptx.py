#!/usr/bin/env python3
"""Convert a PPTX file to a self-contained HTML presentation with formatting preserved."""
import sys
import base64
import io
import re
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("ERROR: python-pptx not installed. Install with: pip install python-pptx")
    sys.exit(1)


def rgb_to_hex(rgb_color):
    """Convert RGBColor to hex string."""
    if rgb_color is None:
        return None
    try:
        return f"#{rgb_color}"
    except:
        return None


def get_text_style(run):
    """Extract inline text styling from a run."""
    styles = []
    try:
        if run.font.bold:
            styles.append("font-weight:bold")
        if run.font.italic:
            styles.append("font-style:italic")
        if run.font.underline:
            styles.append("text-decoration:underline")
        if run.font.size:
            styles.append(f"font-size:{run.font.size.pt}pt")
        if run.font.color and run.font.color.rgb:
            styles.append(f"color:{rgb_to_hex(run.font.color.rgb)}")
        if run.font.name:
            styles.append(f"font-family:'{run.font.name}',sans-serif")
    except:
        pass
    return ";".join(styles)


def get_alignment(paragraph):
    """Get CSS text-align from paragraph alignment."""
    try:
        align = paragraph.alignment
        if align == PP_ALIGN.CENTER:
            return "center"
        elif align == PP_ALIGN.RIGHT:
            return "right"
        elif align == PP_ALIGN.JUSTIFY:
            return "justify"
    except:
        pass
    return "left"


def extract_image(shape):
    """Extract image from shape as base64 data URI."""
    try:
        image = shape.image
        content_type = image.content_type
        image_bytes = image.blob
        b64 = base64.b64encode(image_bytes).decode('utf-8')
        return f"data:{content_type};base64,{b64}"
    except:
        return None


def get_shape_position(shape, slide_width, slide_height):
    """Get shape position as percentages."""
    try:
        left = (shape.left / slide_width) * 100 if shape.left else 0
        top = (shape.top / slide_height) * 100 if shape.top else 0
        width = (shape.width / slide_width) * 100 if shape.width else 50
        height = (shape.height / slide_height) * 100 if shape.height else 30
        return left, top, width, height
    except:
        return 5, 5, 90, 40


def get_slide_background(slide, prs):
    """Extract slide background color from XML."""
    from pptx.oxml.ns import qn
    for source in [slide, slide.slide_layout]:
        try:
            bg_el = source.background._element
            # Look for solidFill > srgbClr inside bgPr
            for sf in bg_el.iter(qn('a:solidFill')):
                clr = sf.find(qn('a:srgbClr'))
                if clr is not None and clr.get('val'):
                    return f"background-color:#{clr.get('val')}"
        except:
            pass
    return "background-color:#ffffff"


def get_shape_fill(shape):
    """Extract shape fill color from XML."""
    from pptx.oxml.ns import qn
    try:
        sp_pr = shape._element.find(qn('p:spPr'))
        if sp_pr is None:
            sp_pr = shape._element.find(qn('a:spPr'))
        if sp_pr is None:
            # Try direct child
            for tag in ['{http://schemas.openxmlformats.org/drawingml/2006/main}spPr',
                        '{http://schemas.openxmlformats.org/presentationml/2006/main}spPr']:
                sp_pr = shape._element.find(tag)
                if sp_pr is not None:
                    break
        if sp_pr is not None:
            sf = sp_pr.find(qn('a:solidFill'))
            if sf is not None:
                clr = sf.find(qn('a:srgbClr'))
                if clr is not None and clr.get('val'):
                    return f"#{clr.get('val')}"
    except:
        pass
    return None


def render_paragraph(paragraph):
    """Render a paragraph with inline formatting."""
    align = get_alignment(paragraph)
    parts = []
    for run in paragraph.runs:
        text = run.text
        if not text:
            continue
        text = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        style = get_text_style(run)
        if style:
            parts.append(f'<span style="{style}">{text}</span>')
        else:
            parts.append(text)
    if not parts:
        return ""
    content = "".join(parts)
    return f'<p style="text-align:{align};margin:0.3em 0;line-height:1.4">{content}</p>'


def convert(pptx_path, output_path=None):
    prs = Presentation(pptx_path)
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    aspect_ratio = slide_width / slide_height if slide_height else 16/9

    slides_html = []

    for i, slide in enumerate(prs.slides, 1):
        bg_style = get_slide_background(slide, prs)
        elements = []

        for shape in sorted(slide.shapes, key=lambda s: (s.top or 0, s.left or 0)):
            left, top, width, height = get_shape_position(shape, slide_width, slide_height)
            pos_style = f"position:absolute;left:{left:.1f}%;top:{top:.1f}%;width:{width:.1f}%;height:{height:.1f}%"

            # Image
            if shape.shape_type == 13 or hasattr(shape, "image"):
                data_uri = extract_image(shape)
                if data_uri:
                    elements.append(
                        f'<div style="{pos_style};display:flex;align-items:center;justify-content:center">'
                        f'<img src="{data_uri}" style="max-width:100%;max-height:100%;object-fit:contain" alt="">'
                        f'</div>'
                    )
                    continue

            # Table
            if shape.has_table:
                table = shape.table
                table_html = '<table style="width:100%;border-collapse:collapse;font-size:0.9em">'
                for row in table.rows:
                    table_html += "<tr>"
                    for cell in row.cells:
                        cell_text = cell.text.replace("&", "&amp;").replace("<", "&lt;")
                        table_html += f'<td style="border:1px solid #ccc;padding:6px 10px">{cell_text}</td>'
                    table_html += "</tr>"
                table_html += "</table>"
                elements.append(f'<div style="{pos_style};overflow:auto">{table_html}</div>')
                continue

            # Text
            if shape.has_text_frame:
                text_parts = []
                for para in shape.text_frame.paragraphs:
                    rendered = render_paragraph(para)
                    if rendered:
                        text_parts.append(rendered)
                if text_parts:
                    content = "".join(text_parts)
                    fill = get_shape_fill(shape)
                    fill_style = f"background-color:{fill};padding:1em;border-radius:8px;" if fill else ""
                    elements.append(
                        f'<div style="{pos_style};{fill_style}overflow:hidden;display:flex;flex-direction:column;justify-content:center">'
                        f'{content}</div>'
                    )
                continue

            # Decorative shape with fill (colored rectangles, bars, etc.)
            fill = get_shape_fill(shape)
            if fill:
                elements.append(
                    f'<div style="{pos_style};background-color:{fill};border-radius:4px"></div>'
                )

        slide_content = "\n".join(elements)
        slides_html.append(
            f'<section class="slide" style="{bg_style}">\n<div class="slide-inner">\n{slide_content}\n</div>\n</section>'
        )

    title = "Presentation"
    # Try to get title from first slide
    if prs.slides:
        for shape in prs.slides[0].shapes:
            if hasattr(shape, "text") and shape.text.strip() and len(shape.text.strip()) < 150:
                title = shape.text.strip()
                break

    html = f'''<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>{title}</title>
<style>
* {{ margin: 0; padding: 0; box-sizing: border-box; }}
html, body {{ height: 100%; overflow: hidden; background: #000; }}
.slide {{
    width: 100vw; height: 100vh;
    display: none;
    align-items: center; justify-content: center;
    overflow: hidden;
}}
.slide.active {{ display: flex; }}
.slide-inner {{
    position: relative;
    width: 1280px; height: {int(1280 / aspect_ratio)}px;
    transform-origin: center center;
    flex-shrink: 0;
}}
.progress {{ position: fixed; bottom: 0; left: 0; height: 4px; background: #0366d6; transition: width 0.3s; z-index: 100; }}
.counter {{ position: fixed; bottom: 12px; right: 20px; font-size: 14px; color: rgba(255,255,255,0.4); z-index: 100; }}
</style>
</head>
<body>
{chr(10).join(slides_html)}
<div class="progress" id="progress"></div>
<div class="counter" id="counter"></div>
<script>
const slides = document.querySelectorAll('.slide');
let current = 0;
function show(n) {{
    slides.forEach(s => s.classList.remove('active'));
    current = Math.max(0, Math.min(n, slides.length - 1));
    slides[current].classList.add('active');
    document.getElementById('progress').style.width = ((current + 1) / slides.length * 100) + '%';
    document.getElementById('counter').textContent = (current + 1) + ' / ' + slides.length;
}}
document.addEventListener('keydown', e => {{
    if (e.key === 'ArrowRight' || e.key === ' ') {{ e.preventDefault(); show(current + 1); }}
    if (e.key === 'ArrowLeft') {{ e.preventDefault(); show(current - 1); }}
}});
let touchStartX = 0;
document.addEventListener('touchstart', e => {{ touchStartX = e.changedTouches[0].screenX; }});
document.addEventListener('touchend', e => {{
    const diff = e.changedTouches[0].screenX - touchStartX;
    if (Math.abs(diff) > 50) {{ diff > 0 ? show(current - 1) : show(current + 1); }}
}});
document.addEventListener('click', e => {{
    if (e.clientX > window.innerWidth / 2) show(current + 1);
    else show(current - 1);
}});
show(0);
function scaleSlides() {{
    document.querySelectorAll('.slide-inner').forEach(inner => {{
        const scaleX = window.innerWidth / inner.offsetWidth;
        const scaleY = window.innerHeight / inner.offsetHeight;
        const scale = Math.min(scaleX, scaleY);
        inner.style.transform = 'scale(' + scale + ')';
    }});
}}
window.addEventListener('resize', scaleSlides);
scaleSlides();
</script>
</body></html>'''

    output = output_path or str(Path(pptx_path).with_suffix('.html'))
    Path(output).write_text(html, encoding='utf-8')
    print(f"Converted to: {output}")
    print(f"Slides: {len(slides_html)}")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: convert-pptx.py <file.pptx> [output.html]")
        sys.exit(1)
    convert(sys.argv[1], sys.argv[2] if len(sys.argv) > 2 else None)
